"""Generate the presentation .pptx with EU-themed design."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from lxml import etree

# EU palette
EU_BLUE = RGBColor(0x00, 0x33, 0x99)
EU_BLUE_DARK = RGBColor(0x00, 0x22, 0x66)
EU_YELLOW = RGBColor(0xFF, 0xCC, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GREY = RGBColor(0xF2, 0xF4, 0xF8)
TEXT = RGBColor(0x1A, 0x1A, 0x2E)
TEXT_LIGHT = RGBColor(0x55, 0x55, 0x66)
ACCENT_TEAL = RGBColor(0x00, 0x99, 0xCC)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
SIDEBAR_W = Inches(2.0)


def _no_outline(shape):
    shape.line.fill.background()


def add_rect(slide, left, top, width, height, fill, *, outline=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if not outline:
        _no_outline(shape)
    return shape


def add_textbox(slide, left, top, width, height, text, *,
                size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, font="Calibri"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb, tf


def add_paragraph(tf, text, *, size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
                  font="Calibri", level=0, space_before=4):
    p = tf.add_paragraph()
    p.alignment = align
    p.level = level
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p


def draw_sidebar(slide, slide_no, total, section, accent=False):
    # Solid blue sidebar
    add_rect(slide, 0, 0, SIDEBAR_W, SLIDE_H, EU_BLUE)

    # Decorative yellow accent at top
    add_rect(slide, 0, Inches(0.6), SIDEBAR_W, Inches(0.05), EU_YELLOW)

    # 12 stars motif (EU flag) - small dots in a circle pattern
    cx, cy = Inches(1.0), Inches(1.7)
    import math
    for i in range(12):
        angle = math.radians(i * 30 - 90)
        r = Inches(0.55)
        x = cx + Emu(int(r * math.cos(angle)))
        y = cy + Emu(int(r * math.sin(angle)))
        size = Inches(0.13)
        star = slide.shapes.add_shape(MSO_SHAPE.STAR_5_POINT,
                                      x - size / 2, y - size / 2, size, size)
        star.fill.solid()
        star.fill.fore_color.rgb = EU_YELLOW
        _no_outline(star)

    # Section label (vertical, near bottom)
    add_textbox(slide, Inches(0.2), Inches(3.0), Inches(1.8), Inches(2.8),
                section, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # Slide number bottom
    add_textbox(slide, Inches(0.2), Inches(6.6), Inches(1.6), Inches(0.5),
                f"{slide_no:02d} / {total:02d}", size=14, bold=True,
                color=EU_YELLOW, align=PP_ALIGN.CENTER)


def draw_header(slide, title, subtitle=None):
    # Title text
    title_left = SIDEBAR_W + Inches(0.5)
    add_textbox(slide, title_left, Inches(0.5), Inches(10.5), Inches(0.9),
                title, size=32, bold=True, color=EU_BLUE_DARK)
    # Yellow underline
    add_rect(slide, title_left, Inches(1.35), Inches(1.2), Inches(0.07), EU_YELLOW)
    if subtitle:
        add_textbox(slide, title_left, Inches(1.5), Inches(10.5), Inches(0.5),
                    subtitle, size=16, color=TEXT_LIGHT)


def add_bullet_list(slide, items, *, top=2.0, left=None, width=10.5, height=5.0,
                    main_size=22, sub_size=18):
    """items: list of (text, level, optional_color)."""
    if left is None:
        left = SIDEBAR_W + Inches(0.5)
    else:
        left = Inches(left)
    tb = slide.shapes.add_textbox(left, Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_top = Emu(0)
    for i, item in enumerate(items):
        text, level = item[0], item[1]
        color = item[2] if len(item) > 2 else (EU_BLUE_DARK if level == 0 else TEXT)
        size = main_size if level == 0 else sub_size
        bold = level == 0
        if i == 0:
            p = tf.paragraphs[0]
            p.text = ""
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.level = 0
        p.space_before = Pt(10 if level == 0 else 4)
        # Custom bullet by prefix
        marker = "▎ " if level == 0 else "    • "
        run = p.add_run()
        run.text = marker
        run.font.name = "Calibri"
        run.font.size = Pt(size)
        run.font.bold = True
        run.font.color.rgb = EU_YELLOW if level == 0 else ACCENT_TEAL
        run2 = p.add_run()
        run2.text = text
        run2.font.name = "Calibri"
        run2.font.size = Pt(size)
        run2.font.bold = bold
        run2.font.color.rgb = color


def make_blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# =====================================================================

def slide_title(prs, total):
    s = make_blank(prs)
    # Full background blue
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, EU_BLUE)
    # Diagonal accent stripes (imitating blockchain blocks)
    for i, x_in in enumerate([8.5, 9.5, 10.5, 11.5]):
        h = Inches(7.5)
        add_rect(s, Inches(x_in), 0, Inches(0.6), h,
                 EU_BLUE_DARK if i % 2 == 0 else EU_BLUE)
    # Big yellow band
    add_rect(s, 0, Inches(2.3), SLIDE_W, Inches(0.08), EU_YELLOW)
    add_rect(s, 0, Inches(5.0), SLIDE_W, Inches(0.04), EU_YELLOW)

    # 12 stars circle (decoration left-top)
    import math
    cx, cy = Inches(1.6), Inches(1.4)
    for i in range(12):
        angle = math.radians(i * 30 - 90)
        r = Inches(0.7)
        x = cx + Emu(int(r * math.cos(angle)))
        y = cy + Emu(int(r * math.sin(angle)))
        size = Inches(0.18)
        st = s.shapes.add_shape(MSO_SHAPE.STAR_5_POINT,
                                x - size / 2, y - size / 2, size, size)
        st.fill.solid(); st.fill.fore_color.rgb = EU_YELLOW; _no_outline(st)

    # Title
    add_textbox(s, Inches(0.6), Inches(2.7), Inches(11.5), Inches(1.0),
                "The Role of Blockchain Infrastructure",
                size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.6), Inches(3.7), Inches(11.5), Inches(0.8),
                "in Shaping the Digital Society of the European Union",
                size=28, bold=False, color=EU_YELLOW, align=PP_ALIGN.CENTER)

    # Author block
    add_textbox(s, Inches(0.6), Inches(5.4), Inches(11.5), Inches(0.6),
                "Totskiy Vjatseslav", size=24, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.6), Inches(5.95), Inches(11.5), Inches(0.5),
                "Group 5151003/40001  •  Peter the Great St. Petersburg Polytechnic University",
                size=16, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.6), Inches(6.45), Inches(11.5), Inches(0.5),
                "Spring 2025–2026",
                size=14, color=EU_YELLOW, align=PP_ALIGN.CENTER)


def slide_outline(prs, total, n):
    s = make_blank(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    draw_sidebar(s, n, total, "Outline")
    draw_header(s, "Agenda", "Five steps to understand blockchain in the EU")

    items = [
        ("What is blockchain & why it matters for the EU", "01"),
        ("EU strategy and regulatory framework", "02"),
        ("Practical applications of blockchain in EU services", "03"),
        ("Challenges and limitations", "04"),
        ("Conclusion", "05"),
    ]
    # Numbered cards
    top = Inches(2.1)
    for i, (text, num) in enumerate(items):
        y = top + Inches(0.85 * i)
        # Number badge
        badge = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                   SIDEBAR_W + Inches(0.5), y, Inches(0.65), Inches(0.65))
        badge.fill.solid(); badge.fill.fore_color.rgb = EU_YELLOW; _no_outline(badge)
        add_textbox(s, SIDEBAR_W + Inches(0.5), y + Inches(0.05),
                    Inches(0.65), Inches(0.55), num,
                    size=18, bold=True, color=EU_BLUE_DARK,
                    align=PP_ALIGN.CENTER)
        # Text
        add_textbox(s, SIDEBAR_W + Inches(1.4), y + Inches(0.07),
                    Inches(9.0), Inches(0.6), text,
                    size=22, bold=True, color=EU_BLUE_DARK)


def slide_intro(prs, total, n):
    s = make_blank(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    draw_sidebar(s, n, total, "Introduction")
    draw_header(s, "Blockchain & the EU Digital Strategy",
                "Why a public ledger matters for a public union")

    # Left column - definition card
    left = SIDEBAR_W + Inches(0.5)
    card = add_rect(s, left, Inches(2.1), Inches(5.0), Inches(4.4), LIGHT_GREY)
    add_textbox(s, left + Inches(0.3), Inches(2.3), Inches(4.5), Inches(0.5),
                "WHAT IS BLOCKCHAIN?", size=14, bold=True, color=EU_BLUE)
    add_textbox(s, left + Inches(0.3), Inches(2.8), Inches(4.5), Inches(0.6),
                "A decentralized digital ledger", size=22, bold=True, color=EU_BLUE_DARK)
    # Three pillars
    pillars = [("Transparency", "Public verifiability"),
               ("Security", "Cryptographic integrity"),
               ("Immutability", "Tamper-proof records")]
    for i, (h, sub) in enumerate(pillars):
        py = Inches(3.7 + i * 0.85)
        bullet = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    left + Inches(0.3), py + Inches(0.1),
                                    Inches(0.25), Inches(0.25))
        bullet.fill.solid(); bullet.fill.fore_color.rgb = EU_YELLOW; _no_outline(bullet)
        add_textbox(s, left + Inches(0.7), py, Inches(4.0), Inches(0.4),
                    h, size=18, bold=True, color=EU_BLUE_DARK)
        add_textbox(s, left + Inches(0.7), py + Inches(0.4), Inches(4.0), Inches(0.4),
                    sub, size=14, color=TEXT_LIGHT)

    # Right column - thesis
    right = SIDEBAR_W + Inches(5.8)
    add_rect(s, right, Inches(2.1), Inches(4.7), Inches(4.4), EU_BLUE)
    add_rect(s, right, Inches(2.1), Inches(4.7), Inches(0.08), EU_YELLOW)
    add_textbox(s, right + Inches(0.3), Inches(2.4), Inches(4.1), Inches(0.5),
                "THESIS", size=14, bold=True, color=EU_YELLOW)
    add_textbox(s, right + Inches(0.3), Inches(2.95), Inches(4.1), Inches(3.4),
                "Blockchain infrastructure plays a key role in building a "
                "secure, transparent, and efficient digital society "
                "in the European Union.",
                size=20, bold=True, color=WHITE)
    # quote-mark decoration
    add_textbox(s, right + Inches(3.7), Inches(5.4), Inches(0.8), Inches(0.8),
                "”", size=72, bold=True, color=EU_YELLOW, align=PP_ALIGN.RIGHT)


def slide_strategy(prs, total, n):
    s = make_blank(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    draw_sidebar(s, n, total, "Strategy")
    draw_header(s, "EU Blockchain Strategy & Regulation",
                "Four pillars of the EU's institutional approach")

    # 4 cards in a 2x2 grid
    cards = [
        ("EBP", "European Blockchain Partnership", "2018",
         "27 EU states + Norway + Liechtenstein  [1]"),
        ("EBSI", "European Blockchain Services Infrastructure", "since 2019",
         "“peer-to-peer network of interconnected nodes running a blockchain-based services infrastructure”  [1 – para. 1]"),
        ("MiCA", "Markets in Crypto-Assets Regulation", "Dec 2024",
         "“MiCA establishes uniform EU market rules for crypto-assets”  [2 – para. 1]"),
        ("EUROPEUM-EDIC", "Governance entity for EBSI", "May 2024",
         "Prepares the infrastructure for full-scale production  [3]"),
    ]
    base_left = SIDEBAR_W + Inches(0.5)
    base_top = Inches(2.1)
    card_w = Inches(5.2)
    card_h = Inches(2.15)
    gap_x = Inches(0.3)
    gap_y = Inches(0.2)
    for i, (tag, title, when, body) in enumerate(cards):
        col, row = i % 2, i // 2
        x = base_left + col * (card_w + gap_x)
        y = base_top + row * (card_h + gap_y)
        # card bg
        add_rect(s, x, y, card_w, card_h, LIGHT_GREY)
        # left tag stripe
        add_rect(s, x, y, Inches(0.18), card_h, EU_BLUE)
        # tag
        add_textbox(s, x + Inches(0.4), y + Inches(0.1), Inches(2.5), Inches(0.5),
                    tag, size=18, bold=True, color=EU_BLUE_DARK)
        # date pill
        add_textbox(s, x + card_w - Inches(1.5), y + Inches(0.1),
                    Inches(1.3), Inches(0.4),
                    when, size=12, bold=True,
                    color=EU_BLUE, align=PP_ALIGN.RIGHT)
        # title
        add_textbox(s, x + Inches(0.4), y + Inches(0.55), card_w - Inches(0.6),
                    Inches(0.5), title, size=14, color=TEXT_LIGHT)
        # body
        add_textbox(s, x + Inches(0.4), y + Inches(1.0), card_w - Inches(0.6),
                    Inches(1.1), body, size=14, color=TEXT)


def slide_identity(prs, total, n):
    s = make_blank(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    draw_sidebar(s, n, total, "Applications")
    draw_header(s, "Application 1: Digital Identity & Diplomas",
                "Self-Sovereign Identity across borders")

    left = SIDEBAR_W + Inches(0.5)
    # Big quote/concept
    add_rect(s, left, Inches(2.1), Inches(10.5), Inches(1.2), EU_BLUE)
    add_textbox(s, left + Inches(0.4), Inches(2.25), Inches(1.5), Inches(0.5),
                "SSI", size=18, bold=True, color=EU_YELLOW)
    add_textbox(s, left + Inches(0.4), Inches(2.65), Inches(10.0), Inches(0.55),
                "Citizens manage their own digital identities — without relying on "
                "centralized authorities  [4]", size=18, color=WHITE)

    # Three step pipeline
    steps = [
        ("1.  Issuance", "University issues a verifiable diploma on EBSI"),
        ("2.  Verification", "Cross-border check via shared infrastructure"),
        ("3.  Recognition", "Credential accepted across all EU states"),
    ]
    sx = left
    sy = Inches(3.7)
    sw = Inches(3.4)
    sh = Inches(2.0)
    gap = Inches(0.2)
    for i, (h, body) in enumerate(steps):
        x = sx + i * (sw + gap)
        add_rect(s, x, sy, sw, sh, LIGHT_GREY)
        add_rect(s, x, sy, sw, Inches(0.08), EU_YELLOW)
        add_textbox(s, x + Inches(0.3), sy + Inches(0.25), sw - Inches(0.4),
                    Inches(0.6), h, size=18, bold=True, color=EU_BLUE_DARK)
        add_textbox(s, x + Inches(0.3), sy + Inches(0.85), sw - Inches(0.4),
                    Inches(1.0), body, size=15, color=TEXT)
        # arrow between steps
        if i < len(steps) - 1:
            ax = x + sw + Inches(0.02)
            ay = sy + Inches(0.85)
            arr = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, ax, ay,
                                     Inches(0.16), Inches(0.3))
            arr.fill.solid(); arr.fill.fore_color.rgb = EU_BLUE; _no_outline(arr)

    # Outcome banner
    add_rect(s, left, Inches(5.95), Inches(10.5), Inches(0.7), EU_YELLOW)
    add_textbox(s, left + Inches(0.4), Inches(6.05), Inches(10.0), Inches(0.5),
                "Result: less bureaucracy, no manual document verification",
                size=18, bold=True, color=EU_BLUE_DARK)


def slide_supply(prs, total, n):
    s = make_blank(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    draw_sidebar(s, n, total, "Applications")
    draw_header(s, "Application 2: Supply Chain Transparency",
                "Immutable records of product origin and movement")

    left = SIDEBAR_W + Inches(0.5)
    # Three big use-case tiles
    tiles = [
        ("Food safety", "Track origin and storage of fresh produce", EU_BLUE),
        ("Pharmaceuticals", "Verify drug provenance and prevent counterfeits", ACCENT_TEAL),
        ("High-value goods", "Authenticate luxury items in cross-border trade", EU_BLUE_DARK),
    ]
    tx = left
    ty = Inches(2.2)
    tw = Inches(3.4)
    th = Inches(3.4)
    gap = Inches(0.2)
    for i, (h, body, color) in enumerate(tiles):
        x = tx + i * (tw + gap)
        add_rect(s, x, ty, tw, th, color)
        # icon-area top half
        add_rect(s, x, ty, tw, Inches(1.4), EU_BLUE_DARK)
        # number
        add_textbox(s, x + Inches(0.3), ty + Inches(0.3), Inches(2.5), Inches(0.7),
                    f"0{i+1}", size=44, bold=True, color=EU_YELLOW)
        # heading
        add_textbox(s, x + Inches(0.3), ty + Inches(1.6), tw - Inches(0.6),
                    Inches(0.5), h, size=20, bold=True, color=WHITE)
        # body
        add_textbox(s, x + Inches(0.3), ty + Inches(2.1), tw - Inches(0.6),
                    Inches(1.2), body, size=14, color=WHITE)

    # Bottom takeaway
    add_textbox(s, left, Inches(5.85), Inches(10.5), Inches(0.6),
                "Builds public trust in the European single market",
                size=18, bold=True, color=EU_BLUE_DARK, align=PP_ALIGN.CENTER)


def slide_challenges(prs, total, n):
    s = make_blank(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    draw_sidebar(s, n, total, "Challenges")
    draw_header(s, "Challenges & Limitations",
                "What still stands in the way of full adoption")

    left = SIDEBAR_W + Inches(0.5)
    # 2x2 grid of warning cards
    items = [
        ("⚙", "Scalability", "Limited transaction throughput across networks"),
        ("⚡", "Energy use", "Mitigated by EBSI's Proof of Authority  [5]"),
        ("⚖", "Regulatory fragmentation", "MiCA transitional measures applied at varying speeds  [6]"),
        ("👥", "Public adoption", "Insufficient understanding of blockchain among citizens"),
    ]
    base_top = Inches(2.1)
    card_w = Inches(5.2)
    card_h = Inches(2.15)
    gap_x = Inches(0.3)
    gap_y = Inches(0.2)
    for i, (icon, title, body) in enumerate(items):
        col, row = i % 2, i // 2
        x = left + col * (card_w + gap_x)
        y = base_top + row * (card_h + gap_y)
        add_rect(s, x, y, card_w, card_h, LIGHT_GREY)
        # icon circle
        circ = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  x + Inches(0.3), y + Inches(0.3),
                                  Inches(1.0), Inches(1.0))
        circ.fill.solid(); circ.fill.fore_color.rgb = EU_BLUE; _no_outline(circ)
        add_textbox(s, x + Inches(0.3), y + Inches(0.35),
                    Inches(1.0), Inches(0.9), icon,
                    size=36, bold=True, color=EU_YELLOW,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # heading
        add_textbox(s, x + Inches(1.5), y + Inches(0.3), card_w - Inches(1.7),
                    Inches(0.6), title, size=18, bold=True, color=EU_BLUE_DARK)
        # body
        add_textbox(s, x + Inches(1.5), y + Inches(0.95), card_w - Inches(1.7),
                    Inches(1.1), body, size=14, color=TEXT)


def slide_conclusion(prs, total, n):
    s = make_blank(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    draw_sidebar(s, n, total, "Conclusion")
    draw_header(s, "Conclusion", "Three takeaways")

    left = SIDEBAR_W + Inches(0.5)
    # Three big takeaways
    items = [
        ("01", "Blockchain reshapes EU digital society",
         "Identity, verification, supply chains"),
        ("02", "EU is fully committed",
         "EBSI, EBP, MiCA, EUROPEUM-EDIC"),
        ("03", "Investment + balanced regulation",
         "Essential for a decentralized digital Europe"),
    ]
    bx = left
    by = Inches(2.2)
    bw = Inches(3.4)
    bh = Inches(3.6)
    gap = Inches(0.2)
    for i, (num, h, body) in enumerate(items):
        x = bx + i * (bw + gap)
        add_rect(s, x, by, bw, bh, WHITE, outline=False)
        # frame
        frame = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, by, bw, bh)
        frame.fill.background()
        frame.line.color.rgb = EU_BLUE
        frame.line.width = Pt(1.5)
        # top ribbon
        add_rect(s, x, by, bw, Inches(0.4), EU_BLUE)
        add_textbox(s, x + Inches(0.3), by + Inches(0.05), Inches(2.0),
                    Inches(0.35), f"Takeaway {num}", size=12, bold=True, color=EU_YELLOW)
        # heading
        add_textbox(s, x + Inches(0.3), by + Inches(0.7), bw - Inches(0.6),
                    Inches(1.4), h, size=20, bold=True, color=EU_BLUE_DARK)
        # body
        add_textbox(s, x + Inches(0.3), by + Inches(2.3), bw - Inches(0.6),
                    Inches(1.2), body, size=15, color=TEXT)

    # Bottom thesis restate
    add_rect(s, left, Inches(6.1), Inches(10.5), Inches(0.6), EU_YELLOW)
    add_textbox(s, left + Inches(0.4), Inches(6.18), Inches(10.0), Inches(0.45),
                "Blockchain is not a buzzword — it is the backbone of a more open EU",
                size=16, bold=True, color=EU_BLUE_DARK)


def slide_references(prs, total, n):
    s = make_blank(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    draw_sidebar(s, n, total, "References")
    draw_header(s, "References")

    refs = [
        ("[1]", "European Commission (2024)",
         "European Blockchain Services Infrastructure"),
        ("[2]", "ESMA (2025)",
         "Markets in Crypto-Assets Regulation (MiCA)"),
        ("[3]", "European Commission (2024)",
         "About EUROPEUM-EDIC"),
        ("[4]", "EQAR (2025)",
         "European Blockchain Service Infrastructure (EBSI)"),
        ("[5]", "MITA (2025)",
         "The European Blockchain Services Infrastructure (EBSI)"),
        ("[6]", "Hogan Lovells (2025)",
         "The EU's Markets in Crypto-Assets MiCA Regulation — A Status Update"),
    ]
    left = SIDEBAR_W + Inches(0.5)
    top = Inches(2.0)
    for i, (num, source, title) in enumerate(refs):
        y = top + Inches(0.78 * i)
        # number badge
        add_rect(s, left, y + Inches(0.1), Inches(0.7), Inches(0.5), EU_BLUE)
        add_textbox(s, left, y + Inches(0.15), Inches(0.7), Inches(0.4),
                    num, size=14, bold=True, color=EU_YELLOW,
                    align=PP_ALIGN.CENTER)
        # source
        add_textbox(s, left + Inches(0.85), y + Inches(0.05),
                    Inches(4.0), Inches(0.4),
                    source, size=15, bold=True, color=EU_BLUE_DARK)
        # title
        add_textbox(s, left + Inches(0.85), y + Inches(0.4),
                    Inches(9.5), Inches(0.4),
                    title, size=13, color=TEXT_LIGHT)


def slide_thanks(prs, total, n):
    s = make_blank(prs)
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, EU_BLUE)
    # Vertical accent stripes
    for i, x_in in enumerate([0.0, 0.4, 0.8]):
        add_rect(s, Inches(x_in), 0, Inches(0.15), SLIDE_H, EU_YELLOW)
    for i, x_in in enumerate([12.55, 12.95, 13.20]):
        add_rect(s, Inches(x_in), 0, Inches(0.12), SLIDE_H, EU_YELLOW)
    # Big stars circle bottom-right corner
    import math
    cx, cy = Inches(11.0), Inches(6.0)
    for i in range(12):
        angle = math.radians(i * 30 - 90)
        r = Inches(0.9)
        x = cx + Emu(int(r * math.cos(angle)))
        y = cy + Emu(int(r * math.sin(angle)))
        size = Inches(0.22)
        st = s.shapes.add_shape(MSO_SHAPE.STAR_5_POINT,
                                x - size / 2, y - size / 2, size, size)
        st.fill.solid(); st.fill.fore_color.rgb = EU_YELLOW; _no_outline(st)

    add_textbox(s, Inches(0.6), Inches(2.5), Inches(12.0), Inches(1.5),
                "Thank You", size=80, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)
    add_rect(s, Inches(5.5), Inches(4.1), Inches(2.3), Inches(0.07), EU_YELLOW)
    add_textbox(s, Inches(0.6), Inches(4.4), Inches(12.0), Inches(0.7),
                "Questions & Discussion", size=28,
                color=EU_YELLOW, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.6), Inches(5.3), Inches(12.0), Inches(0.5),
                "Totskiy Vjatseslav  •  Group 5151003/40001",
                size=16, color=WHITE, align=PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [slide_title, slide_outline, slide_intro, slide_strategy,
                slide_identity, slide_supply, slide_challenges,
                slide_conclusion, slide_references, slide_thanks]
    total = len(builders)
    for i, fn in enumerate(builders, start=1):
        if fn is slide_title:
            fn(prs, total)
        else:
            fn(prs, total, i)

    out = r"c:\Users\vyach\Study-Materials\Semester-4\English\Presentation_Blockchain_EU.pptx"
    prs.save(out)
    print(f"Saved: {out} ({total} slides)")


if __name__ == "__main__":
    main()
