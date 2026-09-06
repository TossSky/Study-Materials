# -*- coding: utf-8 -*-
"""Извлекает картинку-фон каждого слайда из всех pptx в рабочую папку %TEMP%/lec_img.
Для слайда берём картинку максимальной площади (слайды-лекции = полноэкранная картинка).
Пишет manifest.json с инфой по каждому слайду."""
import os, glob, io, json
from pptx import Presentation
from pptx.util import Emu
from PIL import Image

L  = r"c:\Users\vyach\Study-Materials\Semester-4\TeorVer\Lectures"
WORK = os.path.join(os.environ["TEMP"], "lec_img")
os.makedirs(WORK, exist_ok=True)

def slide_emu(prs):
    return prs.slide_width, prs.slide_height

manifest = []
summary = []
for sub in ("ProbabilityTheory", "MathStatistics"):
    d = os.path.join(L, sub)
    for f in sorted(glob.glob(os.path.join(d, "*.pptx"))):
        stem = os.path.splitext(os.path.basename(f))[0]
        prs = Presentation(f)
        sw, sh = slide_emu(prs)
        outdir = os.path.join(WORK, sub, stem)
        os.makedirs(outdir, exist_ok=True)
        no_pic = 0
        for i, slide in enumerate(prs.slides, 1):
            # найдём картинку максимальной площади
            best = None; best_area = -1
            n_pics = 0; native_text = ""
            for sh_ in slide.shapes:
                if sh_.shape_type == 13:  # PICTURE
                    n_pics += 1
                    try:
                        area = (sh_.width or 0) * (sh_.height or 0)
                    except Exception:
                        area = 0
                    if area > best_area:
                        best_area = area; best = sh_
                if sh_.has_text_frame and sh_.text_frame.text.strip():
                    native_text += sh_.text_frame.text.strip() + "\n"
            png = os.path.join(outdir, f"slide{i:02d}.png")
            saved = False; w=h=0
            if best is not None:
                try:
                    blob = best.image.blob
                    im = Image.open(io.BytesIO(blob)).convert("RGB")
                    im.save(png, "PNG")
                    w, h = im.size
                    saved = True
                except Exception as e:
                    pass
            if not saved:
                no_pic += 1
            manifest.append({
                "sub": sub, "lecture": stem, "slide": i,
                "png": png if saved else None,
                "img_w": w, "img_h": h,
                "n_pics": n_pics,
                "native_text": native_text.strip(),
            })
        summary.append(f"{sub}/{stem}: слайдов={len(prs.slides)} без_картинки={no_pic}")

with open(os.path.join(WORK, "manifest.json"), "w", encoding="utf-8") as fp:
    json.dump(manifest, fp, ensure_ascii=False, indent=1)

print("WORK =", WORK)
print("всего записей:", len(manifest))
nopic = [m for m in manifest if not m["png"]]
print("слайдов без извлечённой картинки:", len(nopic))
print("--- сводка по лекциям (первые/последние) ---")
for s in summary:
    print(" ", s)
