"""
Генератор презентации к защите 1 пункта курсовой по ОС.
Тема: «Изоляция Ethereum-ноды средствами безопасности ядра Linux».
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]

# --------- Палитра ---------
C_DARK = RGBColor(0x14, 0x1B, 0x2D)
C_ACCENT = RGBColor(0xF7, 0x93, 0x1E)
C_TEXT = RGBColor(0xE4, 0xE6, 0xEB)
C_MUTED = RGBColor(0x9A, 0xA5, 0xB8)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, *, size=20, bold=False,
             color=C_TEXT, align=PP_ALIGN.LEFT, font='Calibri'):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split('\n') if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_bullets(slide, left, top, width, height, items, *, size=20, color=C_TEXT, bullet='—'):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = f'{bullet}  {item}'
        run.font.name = 'Calibri'
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return box


def add_accent_bar(slide, left, top, width=0.3, height=0.08, color=C_ACCENT):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(left), Inches(top),
                                 Inches(width), Inches(height))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp


def add_footer(slide, num, total):
    add_text(slide, 0.5, 7.1, 10, 0.3,
             'Курсовая работа по ОС · Тоцкий В. · гр. 5151003/40001',
             size=10, color=C_MUTED)
    add_text(slide, 12.2, 7.1, 0.8, 0.3,
             f'{num} / {total}', size=10, color=C_MUTED, align=PP_ALIGN.RIGHT)


TOTAL = 12


# ============================================================
# Слайд 1 — Титульный
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, C_DARK)

add_accent_bar(s, 0.6, 2.3, width=1.5, height=0.1)

add_text(s, 0.6, 2.5, 12, 0.5,
         'КУРСОВАЯ РАБОТА ПО ДИСЦИПЛИНЕ «ОПЕРАЦИОННЫЕ СИСТЕМЫ»',
         size=14, color=C_ACCENT, bold=True)

add_text(s, 0.6, 3.1, 12, 1.5,
         'Изоляция Ethereum-ноды\nсредствами безопасности ядра Linux',
         size=40, bold=True, color=C_WHITE)

add_text(s, 0.6, 5.2, 12, 0.4,
         'Защита промежуточного варианта (1 пункт)',
         size=20, color=C_TEXT)

add_text(s, 0.6, 6.1, 6, 0.3,
         'Выполнил: студент гр. 5151003/40001 В. Тоцкий',
         size=14, color=C_MUTED)
add_text(s, 0.6, 6.4, 6, 0.3,
         'Руководитель: ст. преп. Г.Д. Гавва',
         size=14, color=C_MUTED)
add_text(s, 0.6, 6.7, 6, 0.3,
         'СПбПУ, ИКНК, ВШКБ · 4 семестр · 2025/2026',
         size=14, color=C_MUTED)


# ============================================================
# Слайд 2 — Актуальность
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, C_DARK)
add_accent_bar(s, 0.6, 0.6)
add_text(s, 0.6, 0.8, 12, 0.6, 'Актуальность', size=28, bold=True, color=C_WHITE)

add_bullets(s, 0.6, 1.8, 12, 5.0, [
    'Ethereum-нода — инфраструктурный компонент Web3, обрабатывающий недоверенный сетевой трафик',
    'Процесс хранит криптографические ключи и подписывает транзакции; компрометация → прямое хищение активов',
    'Большая кодовая база клиентов (Geth, Reth, Erigon) и нативные зависимости → потенциальные уязвимости парсеров',
    'Линия карьерных интересов: Web3 Security / Smart Contract Auditing (SingularityNET, ASI Chain)',
    'Linux предоставляет штатные механизмы изоляции — seccomp, namespaces, cgroups — без внешних средств защиты',
    'Задача изоляции процесса на уровне ядра ОС соответствует профилю дисциплины «Операционные системы»',
], size=18)
add_footer(s, 2, TOTAL)


# ============================================================
# Слайд 3 — Цель и задачи
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, C_DARK)
add_accent_bar(s, 0.6, 0.6)
add_text(s, 0.6, 0.8, 12, 0.6, 'Цель и задачи', size=28, bold=True, color=C_WHITE)

add_text(s, 0.6, 1.7, 12, 0.5, 'Цель', size=18, bold=True, color=C_ACCENT)
add_text(s, 0.6, 2.2, 12, 0.9,
         'Применить штатные механизмы ядра Linux для защиты процесса Ethereum-ноды: '
         'минимизировать поверхность атаки и ограничить права процесса на уровне ядра ОС.',
         size=16, color=C_TEXT)

add_text(s, 0.6, 3.4, 12, 0.5, 'Задачи', size=18, bold=True, color=C_ACCENT)
add_bullets(s, 0.6, 3.9, 12, 3.5, [
    'Построить модель угроз для Ethereum-ноды как пользовательского процесса Linux',
    'Изучить механизмы seccomp BPF, namespaces, cgroups v2 — применимость к задаче изоляции',
    'Профилировать системные вызовы ноды Anvil (strace, perf), составить whitelist',
    'Разработать seccomp-профиль и построить изолированную среду (net namespace + cgroup)',
    'Сравнить модель угроз до и после применения защитных мер, оформить стенд',
], size=16)
add_footer(s, 3, TOTAL)


# ============================================================
# Слайд 4 — Объект: Ethereum-нода Anvil
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, C_DARK)
add_accent_bar(s, 0.6, 0.6)
add_text(s, 0.6, 0.8, 12, 0.6, 'Объект: Ethereum-нода (Anvil)', size=28, bold=True, color=C_WHITE)

add_text(s, 0.6, 1.8, 12, 0.4,
         'Anvil — локальная JSON-RPC-нода из набора Foundry, без синхронизации публичной сети',
         size=16, color=C_MUTED)

add_text(s, 0.6, 2.5, 6, 0.5, 'Ядро ОС → набор примитивов', size=18, bold=True, color=C_ACCENT)
add_bullets(s, 0.6, 3.0, 6, 3.5, [
    'epoll_wait, read, write — event loop',
    'sendto, recvfrom — p2p и JSON-RPC сокеты',
    'mmap, pread, pwrite — БД состояния (LevelDB/RocksDB)',
    'clone, futex — многопоточность',
    'openat, fstat, close — файлы и ключи',
], size=14)

add_text(s, 6.8, 2.5, 6, 0.5, 'Особенности процесса', size=18, bold=True, color=C_ACCENT)
add_bullets(s, 6.8, 3.0, 6, 3.5, [
    'Постоянно слушает недоверенный трафик',
    'Работа с приватными ключами в адресном пространстве',
    'Интенсивный I/O по БД состояния',
    'Тяжёлое CPU/RAM при валидации блоков',
    'Поведение эквивалентно промышленным клиентам',
], size=14)
add_footer(s, 4, TOTAL)


# ============================================================
# Слайд 5 — Модель угроз
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, C_DARK)
add_accent_bar(s, 0.6, 0.6)
add_text(s, 0.6, 0.8, 12, 0.6, 'Модель угроз', size=28, bold=True, color=C_WHITE)

add_text(s, 0.6, 1.7, 12, 0.5,
         'Базовое предположение: атакующий получил RCE в контексте процесса ноды, без прав root',
         size=16, color=C_MUTED)

threats = [
    ('RCE',              'Удалённое исполнение кода через парсеры JSON-RPC / p2p'),
    ('LPE',              'Повышение привилегий до суперпользователя через уязвимости ядра'),
    ('Эксфильтрация',    'Чтение /etc/shadow, ключей других процессов, конфигураций'),
    ('DoS',              'Исчерпание CPU, памяти, дескрипторов — удар по системе в целом'),
    ('Сетевые атаки',    'Открытие туннеля наружу, сканирование локальной сети'),
]
for i, (name, desc) in enumerate(threats):
    y = 2.4 + i * 0.8
    add_text(s, 0.6, y, 2.6, 0.6, name, size=18, bold=True, color=C_ACCENT)
    add_text(s, 3.3, y, 9.5, 0.6, desc, size=16, color=C_TEXT)
add_footer(s, 5, TOTAL)


# ============================================================
# Слайд 6 — Три уровня защиты
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, C_DARK)
add_accent_bar(s, 0.6, 0.6)
add_text(s, 0.6, 0.8, 12, 0.6, 'Три уровня защиты средствами ядра Linux',
         size=28, bold=True, color=C_WHITE)

from pptx.enum.shapes import MSO_SHAPE

cols = [
    ('seccomp BPF',    'Фильтр системных вызовов',
     'Whitelist минимально\nнеобходимых syscalls.\nВсё остальное блокируется\nядром (KILL / ERRNO)'),
    ('namespaces',     'Разделение ресурсов',
     'net namespace: изоляция\nсетевых интерфейсов.\nmount namespace: ограничение\nвидимости ФС'),
    ('cgroups v2',     'Квотирование',
     'memory.max, cpu.max,\npids.max, io.max.\nЗащита от DoS и fork-bomb\nиз скомпрометированного процесса'),
]
for i, (title, subtitle, body) in enumerate(cols):
    x = 0.6 + i * 4.2
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(x), Inches(1.8),
                             Inches(4.0), Inches(4.6))
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor(0x1E, 0x2A, 0x44)
    shp.line.color.rgb = C_ACCENT
    shp.line.width = Pt(1.5)

    add_text(s, x + 0.3, 2.0, 3.6, 0.5, title, size=22, bold=True, color=C_ACCENT)
    add_text(s, x + 0.3, 2.6, 3.6, 0.4, subtitle, size=14, color=C_MUTED)
    add_text(s, x + 0.3, 3.2, 3.6, 3.0, body, size=14, color=C_TEXT)
add_footer(s, 6, TOTAL)


# ============================================================
# Слайд 7 — seccomp BPF
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, C_DARK)
add_accent_bar(s, 0.6, 0.6)
add_text(s, 0.6, 0.8, 12, 0.6, 'seccomp BPF: фильтрация системных вызовов',
         size=28, bold=True, color=C_WHITE)

add_bullets(s, 0.6, 1.8, 12, 5.0, [
    'Ядро применяет программу cBPF к каждому syscall процесса',
    'Действия: ALLOW · KILL_PROCESS · ERRNO · TRAP · LOG · USER_NOTIF',
    'Установка: prctl(PR_SET_SECCOMP) или seccomp(2), обязателен no_new_privs',
    'Инструментарий: libseccomp (высокоуровневый API), scmp_sys_resolver',
    'Фильтр наследуется при fork/clone и не может быть ослаблен (монотонно сужает)',
    'Whitelist для ноды ≈ 40–60 syscalls; блокируются mount, init_module, bpf и т. п.',
], size=18)
add_footer(s, 7, TOTAL)


# ============================================================
# Слайд 8 — namespaces
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, C_DARK)
add_accent_bar(s, 0.6, 0.6)
add_text(s, 0.6, 0.8, 12, 0.6, 'namespaces: разделение ресурсов',
         size=28, bold=True, color=C_WHITE)

add_text(s, 0.6, 1.8, 6.2, 0.5, '8 типов пространств имён', size=18, bold=True, color=C_ACCENT)
add_bullets(s, 0.6, 2.3, 6.2, 4.5, [
    'mnt — точки монтирования ФС',
    'pid — иерархия PID',
    'net — сетевые интерфейсы и сокеты',
    'ipc — System V IPC, POSIX-очереди',
    'uts — имя хоста',
    'user — UID/GID',
    'cgroup, time',
], size=15)

add_text(s, 7.1, 1.8, 5.8, 0.5, 'Применение в работе', size=18, bold=True, color=C_ACCENT)
add_bullets(s, 7.1, 2.3, 5.8, 4.5, [
    'net ns — единственный доступный интерфейс: veth',
    'Фиксированный маршрут на хостовый мост',
    'Блокировка открытия произвольных соединений',
    'mount ns — «чистое» представление ФС',
    'Доступны: рабочий каталог ноды, ключи, libc',
    'API: clone(CLONE_NEW*), unshare, setns',
], size=15)
add_footer(s, 8, TOTAL)


# ============================================================
# Слайд 9 — cgroups v2
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, C_DARK)
add_accent_bar(s, 0.6, 0.6)
add_text(s, 0.6, 0.8, 12, 0.6, 'cgroups v2: квотирование ресурсов',
         size=28, bold=True, color=C_WHITE)

add_text(s, 0.6, 1.7, 12, 0.4,
         'Единая иерархия в /sys/fs/cgroup; процесс ↔ cgroup через cgroup.procs',
         size=16, color=C_MUTED)

cgroups = [
    ('memory.max',  'Лимит RAM; OOM-killer завершает только процесс ноды'),
    ('cpu.max',     'Квота CPU в формате «<квота>/<период>» (80000/100000 = 80% ядра)'),
    ('pids.max',    'Максимум процессов — защита от fork bomb'),
    ('io.max',      'Пропускная способность дискового I/O'),
]
for i, (name, desc) in enumerate(cgroups):
    y = 2.4 + i * 0.9
    add_text(s, 0.6, y, 2.8, 0.6, name, size=18, bold=True, color=C_ACCENT, font='Consolas')
    add_text(s, 3.5, y, 9.3, 0.6, desc, size=16, color=C_TEXT)
add_footer(s, 9, TOTAL)


# ============================================================
# Слайд 10 — Практическая часть (план)
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, C_DARK)
add_accent_bar(s, 0.6, 0.6)
add_text(s, 0.6, 0.8, 12, 0.6, 'Практическая часть (разделы 2–3)',
         size=28, bold=True, color=C_WHITE)

steps = [
    'Развернуть Anvil и проверить JSON-RPC',
    'Снять syscall-профиль через strace и perf',
    'Построить whitelist seccomp через libseccomp',
    'Создать net namespace с veth ↔ мост',
    'Завести cgroup с memory/cpu/pids лимитами',
    'Сценарий симулируемой атаки: RCE → эскалация',
    'Сравнить поведение: без защиты vs с изоляцией',
    'Оформить воспроизводимый стенд на GitHub',
]
for i, step in enumerate(steps):
    col = i % 2
    row = i // 2
    x = 0.6 + col * 6.4
    y = 1.8 + row * 1.15
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL,
                             Inches(x), Inches(y),
                             Inches(0.6), Inches(0.6))
    shp.fill.solid()
    shp.fill.fore_color.rgb = C_ACCENT
    shp.line.fill.background()
    add_text(s, x, y + 0.08, 0.6, 0.5, str(i + 1),
             size=18, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    add_text(s, x + 0.9, y + 0.1, 5.3, 0.8, step, size=16, color=C_TEXT)
add_footer(s, 10, TOTAL)


# ============================================================
# Слайд 11 — Ожидаемые результаты
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, C_DARK)
add_accent_bar(s, 0.6, 0.6)
add_text(s, 0.6, 0.8, 12, 0.6, 'Ожидаемые результаты', size=28, bold=True, color=C_WHITE)

add_bullets(s, 0.6, 1.9, 12, 5.0, [
    'Обоснованный seccomp-профиль для Ethereum-ноды (whitelist ≈ 40–60 syscalls)',
    'Работающая изолированная среда запуска: net namespace + cgroup + seccomp',
    'Сравнительный анализ модели угроз до и после применения защитных мер',
    'Документированный воспроизводимый стенд (README + скрипты + конфиги) на GitHub',
    'Опыт, применимый в задачах Web3-инфраструктуры: защита промышленных клиентов (Geth, Reth)',
], size=20)
add_footer(s, 11, TOTAL)


# ============================================================
# Слайд 12 — Спасибо / Вопросы
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, C_DARK)
add_accent_bar(s, 6.3, 3.3, width=0.8, height=0.12)
add_text(s, 0.6, 3.5, 12, 1.0, 'Спасибо за внимание',
         size=48, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(s, 0.6, 4.6, 12, 0.6, 'Готов ответить на вопросы',
         size=22, color=C_ACCENT, align=PP_ALIGN.CENTER)


# ============================================================
# Сохранение
# ============================================================
out_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                        'Презентация_ОС_Тоцкий.pptx')
prs.save(out_path)
print(f'Saved: {out_path}')
